#!/usr/bin/env python3
"""
HKTVmall 商戶招商 PPT 生成器 v2
Flask API 版本 - 生成精美設計的 PPT

使用方法:
  python3 ppt_server.py                    # 啟動 API 服務 (默認 0.0.0.0:5000)
  python3 ppt_server.py --port 8080        # 指定端口
  python3 ppt_server.py --host 192.168.x.x  # 指定主機

API 端點:
  POST /generate
  {
    "merchant_name": "ABC電器",
    "category": "本地"|"海外"|"服務",
    "pain_points": ["logistics", "fee"],
    "contact_email": "aog.merc@hktv.com.hk"
  }
"""

import io
import os
import sys
import argparse
from pathlib import Path

from flask import Flask, request, send_file, jsonify, send_from_directory
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn, nsmap
from pptx.oxml import parse_xml
from lxml import etree
import copy

app = Flask(__name__)

# ============================================================
# 顏色配置（HKTV 主題）
# ============================================================
COLORS = {
    'hkred': RGBColor(230, 0, 18),        # HKTV 主紅
    'gold': RGBColor(255, 184, 0),         # 金色點綴
    'dark': RGBColor(34, 34, 34),         # 深色文字
    'muted': RGBColor(102, 102, 102),    # 灰色副文字
    'white': RGBColor(255, 255, 255),     # 白色
    'accent_blue': RGBColor(0, 102, 204), # 藍色強調
    'green': RGBColor(0, 136, 0),         # 綠色（成功/解決方案）
}

# ============================================================
# 痛點與解決方案
# ============================================================
PAIN_SOLUTIONS = {
    "logistics": {
        "zh": "物流成本高",
        "solution": "全港最大住宅配送網絡，350+輛貨車，500+取貨點，大幅降低物流成本"
    },
    "fee": {
        "zh": "平台費用不透明",
        "solution": "佣金按分類碼固定比例計算，每月15個工作日結算，無隱藏收費"
    },
    "experience": {
        "zh": "缺乏電商經驗",
        "solution": "電商學院提供完整培訓，MMS後台操作、產品上架、廣告投放一對一支援"
    },
    "inventory": {
        "zh": "擔心庫存風險",
        "solution": "GREEN LAB臨期百貨計劃，3PL靈活存貨管理，代售模式先售後買降低風險"
    },
    "brand": {
        "zh": "希望提升品牌知名度",
        "solution": "每年過億廣告投放，TVB黃金時段，58個MTR站，23列全車身，260萬+買家"
    },
    "young": {
        "zh": "希望拓展年輕客群",
        "solution": "精準CRM定位，25-40歲核心消費群，每季度4.7x購買頻率的忠實客戶"
    },
    "exposure": {
        "zh": "希望增加線上曝光",
        "solution": "關鍵字廣告、橫幅廣告、首頁推廣位，SEO及站內搜尋排名優化"
    }
}

# ============================================================
# 招商計劃配置
# ============================================================
PLANS = {
    "本地": {
        "name": "常規商戶加盟計劃 (Type A Plan)",
        "fee": "HK$25,000 + HK$22,000 廣告金",
        "features": [
            "本地常規商戶",
            "3-5個工作天成功開店",
            "商戶管理平台全面掌控",
            "全渠道推廣工具",
            "專業客戶服務支援"
        ]
    },
    "海外": {
        "name": "海外商戶加盟計劃 (Overseas Plan)",
        "fee": "HK$3,000（原 HK$25,000）",
        "features": [
            "海外跨境商戶專屬",
            "拓展國際市場客戶",
            "高效無縫配送解決方案",
            "多語言客戶服務支援"
        ]
    },
    "服務": {
        "name": "服務類商戶加盟計劃 (Service Deal)",
        "fee": "HK$3,000",
        "features": [
            "e-Voucher 核銷系統",
            "O2O 門店網絡",
            "直播頻道推廣機會",
            "會員引流支援"
        ]
    }
}

# ============================================================
# 輔助函數
# ============================================================

def add_text_box(slide, text, left, top, width, height, font_size=14, bold=False,
                color=None, align=PP_ALIGN.LEFT, bg_color=None):
    """添加文本框"""
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.alignment = align
    if color:
        p.font.color.rgb = color
    if bg_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
    return shape


def add_rectangle(slide, left, top, width, height, fill_color=None, line_color=None):
    """添加矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def set_slide_background(slide, color):
    """設置幻燈片背景"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


# ============================================================
# PPT 生成函數
# ============================================================

def create_pptx(merchant_name, category="本地", pain_points=None, contact_email="aog.merc@hktv.com.hk"):
    """生成精美設計的 HKTVmall 招商 PPT"""
    if pain_points is None:
        pain_points = []

    prs = Presentation()
    prs.layout = prs.slide_layouts[6]  # 空白佈局
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    C = COLORS
    plan_info = PLANS.get(category, PLANS["本地"])

    # ────────────────────────────────────────────────────────
    # SLIDE 1: 封面（紅色背景）
    # ────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, C['hkred'])

    # 頂部分隔線
    add_rectangle(slide, 0, 0, 13.333, 0.15, C['gold'])

    # 主標題
    add_text_box(slide, '香港網上購物平台', 0.5, 1.5, 12.333, 1,
                font_size=54, bold=True, color=C['white'], align=PP_ALIGN.CENTER)

    # No.1
    add_text_box(slide, 'No.1', 0.5, 2.8, 12.333, 0.9,
                font_size=48, bold=True, color=C['gold'], align=PP_ALIGN.CENTER)

    # 商戶合作方案
    add_text_box(slide, '商戶合作方案', 0.5, 3.7, 12.333, 0.8,
                font_size=40, color=C['white'], align=PP_ALIGN.CENTER)

    # 商戶名稱
    add_text_box(slide, f'{merchant_name} 專屬方案', 0.5, 5.2, 12.333, 0.6,
                font_size=28, bold=True, color=C['gold'], align=PP_ALIGN.CENTER)

    # 底部裝飾線
    add_rectangle(slide, 0, 7.35, 13.333, 0.15, C['gold'])

    # ────────────────────────────────────────────────────────
    # SLIDE 2: No.1 香港網上購物平台
    # ────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, C['white'])

    # 頂部紅線
    add_rectangle(slide, 0, 0, 13.333, 0.08, C['hkred'])

    add_text_box(slide, 'No.1', 0.5, 0.5, 12.333, 1,
                font_size=64, bold=True, color=C['hkred'], align=PP_ALIGN.CENTER)
    add_text_box(slide, '香港網上購物平台', 0.5, 1.5, 12.333, 1,
                font_size=48, bold=True, color=C['dark'], align=PP_ALIGN.CENTER)
    add_text_box(slide, '24小時網上購物  ·  物流 & O2O 門市  ·  營銷與廣告\n'
                '平台大數據分析  ·  發展海外市場  ·  商客互聯平台與直播頻道',
                0.5, 3, 12.333, 2,
                font_size=24, color=C['dark'], align=PP_ALIGN.CENTER)

    # ────────────────────────────────────────────────────────
    # SLIDE 3: HKTVmall 里程碑
    # ────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, C['white'])
    add_rectangle(slide, 0, 0, 13.333, 0.08, C['hkred'])

    add_text_box(slide, 'HKTVmall 的里程碑', 0.5, 0.3, 12.333, 0.7,
                font_size=32, bold=True, color=C['dark'])
    add_text_box(slide, '城市電訊在香港聯合交易所上市（香港聯交所代號：1137）',
                0.5, 0.9, 12.333, 0.4,
                font_size=12, color=C['muted'], align=PP_ALIGN.CENTER)

    milestones = [
        ('1992', '城市電訊（香港）有限公司在香港正式成立'),
        ('1997', '香港電視面世'),
        ('2000', '海外業務擴張及開拓業務新板塊'),
        ('2012', '電商正式啟動'),
        ('2015-2016', '香港電視正式轉型為網上購物'),
        ('2017-2018', '物流中心自動化揀貨及倉存系統全面運作'),
        ('2019-2020', '海外業務擴張及開拓業務新板塊'),
        ('2021-2022', '街市即日餸推出網路服務、直播頻道'),
        ('2023-現在', '持續引領香港電商市場')
    ]

    y = 1.5
    for year, desc in milestones:
        add_text_box(slide, year, 0.8, y, 1.5, 0.4,
                    font_size=14, bold=True, color=C['hkred'])
        add_text_box(slide, desc, 2.5, y, 10, 0.4,
                    font_size=14, color=C['dark'])
        y += 0.55

    # ────────────────────────────────────────────────────────
    # SLIDE 4-6: GMV 數據
    # ────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, C['white'])
    add_rectangle(slide, 0, 0, 13.333, 0.08, C['hkred'])

    add_text_box(slide, '強大流量與卓越表現', 0.5, 0.3, 12.333, 0.7,
                font_size=32, bold=True, color=C['dark'])
    add_text_box(slide, '香港領先的電子商務平台',
                0.5, 1.2, 12.333, 0.5,
                font_size=20, color=C['muted'], align=PP_ALIGN.CENTER)
    add_text_box(slide, '2025 Total GMV',
                0.5, 2, 12.333, 0.6,
                font_size=24, color=C['muted'], align=PP_ALIGN.CENTER)
    add_text_box(slide, 'HK$7.89 Billion',
                0.5, 2.6, 12.333, 1.5,
                font_size=72, bold=True, color=C['hkred'], align=PP_ALIGN.CENTER)
    add_text_box(slide, '*香港科技探索有限公司2025年年報數據',
                0.5, 4.2, 12.333, 0.4,
                font_size=11, color=C['muted'], align=PP_ALIGN.CENTER)

    # ────────────────────────────────────────────────────────
    # SLIDE 7-10: 物流系統
    # ────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, C['white'])
    add_rectangle(slide, 0, 0, 13.333, 0.08, C['hkred'])

    add_text_box(slide, '物流核心支柱', 0.5, 0.3, 12.333, 0.7,
                font_size=32, bold=True, color=C['dark'])

    logistics_data = [
        ('350+ 輛貨車', '全港最大住宅配送網絡'),
        ('75家 O2O 門店', '提升消費者的線下體驗'),
        ('每日10萬+ 訂單量', '處理每日的最後一哩配送'),
        ('引入德國自動化執貨系統', '機械化自動執貨及倉存系統')
    ]

    y = 1.3
    for num, desc in logistics_data:
        # 數字
        add_text_box(slide, num, 0.8, y, 5, 0.5,
                    font_size=24, bold=True, color=C['hkred'])
        # 描述
        add_text_box(slide, desc, 0.8, y + 0.45, 11, 0.4,
                    font_size=16, color=C['dark'])
        y += 1.1

    # ────────────────────────────────────────────────────────
    # SLIDE 11: 500+ 取貨點
    # ────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, C['white'])
    add_rectangle(slide, 0, 0, 13.333, 0.08, C['hkred'])

    add_text_box(slide, '全港最大住宅配送網絡', 0.5, 0.3, 12.333, 0.7,
                font_size=32, bold=True, color=C['dark'])
    add_text_box(slide, '500+', 0.5, 1.5, 12.333, 2,
                font_size=120, bold=True, color=C['hkred'], align=PP_ALIGN.CENTER)
    add_text_box(slide, '取貨點', 0.5, 3.5, 12.333, 0.8,
                font_size=36, color=C['dark'], align=PP_ALIGN.CENTER)
    add_text_box(slide, '(香港郵政、Circle K 便利店、759阿信屋，真的城等)',
                0.5, 4.5, 12.333, 0.5,
                font_size=14, color=C['muted'], align=PP_ALIGN.CENTER)

    # ────────────────────────────────────────────────────────
    # SLIDE 12: 大型推廣活動
    # ────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, C['white'])
    add_rectangle(slide, 0, 0, 13.333, 0.08, C['hkred'])

    add_text_box(slide, '大型優惠推廣活動', 0.5, 0.3, 12.333, 0.7,
                font_size=32, bold=True, color=C['dark'])
    add_text_box(slide, '4,000萬', 0.5, 1.3, 12.333, 1.5,
                font_size=80, bold=True, color=C['hkred'], align=PP_ALIGN.CENTER)
    add_text_box(slide, '單日銷售額突破', 0.5, 2.8, 12.333, 0.6,
                font_size=24, color=C['dark'], align=PP_ALIGN.CENTER)
    add_text_box(slide, 'HKTVmall全額承擔85折  ·  媒體新聞報導  ·  投放TVB電視廣告',
                0.5, 3.5, 12.333, 0.5,
                font_size=16, color=C['muted'], align=PP_ALIGN.CENTER)

    # ────────────────────────────────────────────────────────
    # SLIDE 13: 蘇寧直播案例
    # ────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, C['white'])
    add_rectangle(slide, 0, 0, 13.333, 0.08, C['hkred'])

    add_text_box(slide, '2025年 蘇寧 x HKTVmall 專屬直播', 0.5, 0.3, 12.333, 0.7,
                font_size=28, bold=True, color=C['dark'])
    add_text_box(slide, '90萬', 0.5, 1.2, 12.333, 1.5,
                font_size=96, bold=True, color=C['hkred'], align=PP_ALIGN.CENTER)
    add_text_box(slide, '一晚GMV接近', 0.5, 2.7, 12.333, 0.5,
                font_size=24, color=C['dark'], align=PP_ALIGN.CENTER)

    case_points = [
        '• 產品數量：1700件+',
        '• 積極參與HKTVmall市場推廣活動',
        '• 季節性曝光 / 預售獨家產品',
        '• HKTVmall獨家最低價格促銷'
    ]
    y = 3.5
    for point in case_points:
        add_text_box(slide, point, 1.5, y, 10, 0.45,
                    font_size=15, color=C['dark'])
        y += 0.5

    # ────────────────────────────────────────────────────────
    # SLIDE 14: 商戶成功案例
    # ────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, C['white'])
    add_rectangle(slide, 0, 0, 13.333, 0.08, C['hkred'])

    add_text_box(slide, '真實商家成效', 0.5, 0.3, 12.333, 0.7,
                font_size=32, bold=True, color=C['dark'])
    add_text_box(slide, '投資回報見證', 0.5, 1, 12.333, 0.5,
                font_size=18, color=C['muted'])

    # 案例卡片
    cases = [
        ('477杯', '中小企新乳酪品牌\n單場直播售出'),
        ('3萬', '中小企新乳酪品牌\n帶來超過生意額'),
        ('1000+件', '2025年入駐新品牌\n第一季已售出'),
        ('3倍', '小品牌\n年增長率')
    ]
    x = 0.8
    for num, desc in cases:
        add_rectangle(slide, x, 1.8, 2.8, 2.5, C['hkred'])
        add_text_box(slide, num, x, 2, 2.8, 0.8,
                    font_size=28, bold=True, color=C['white'], align=PP_ALIGN.CENTER)
        add_text_box(slide, desc, x, 2.9, 2.8, 1,
                    font_size=12, color=C['white'], align=PP_ALIGN.CENTER)
        x += 3.1

    # ────────────────────────────────────────────────────────
    # SLIDE 15: 招商計劃
    # ────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, C['white'])
    add_rectangle(slide, 0, 0, 13.333, 0.08, C['hkred'])

    add_text_box(slide, plan_info['name'], 0.5, 0.3, 12.333, 0.7,
                font_size=28, bold=True, color=C['dark'])

    # 年費
    add_rectangle(slide, 0.5, 1.2, 4, 1.2, C['hkred'])
    add_text_box(slide, '開店費用', 0.7, 1.3, 3.5, 0.4,
                font_size=14, color=C['white'])
    add_text_box(slide, plan_info['fee'], 0.7, 1.7, 3.5, 0.5,
                font_size=18, bold=True, color=C['white'])

    # 佣金
    add_text_box(slide, '銷售佣金', 5, 1.3, 4, 0.4,
                font_size=14, color=C['muted'])
    add_text_box(slide, '根據產品類別計算', 5, 1.7, 4, 0.5,
                font_size=16, color=C['dark'])

    # 特點
    y = 2.8
    for feature in plan_info['features']:
        add_text_box(slide, f'✓  {feature}', 0.8, y, 11, 0.5,
                    font_size=18, color=C['dark'])
        y += 0.6

    # ────────────────────────────────────────────────────────
    # SLIDE 16: 簡易加入流程
    # ────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, C['white'])
    add_rectangle(slide, 0, 0, 13.333, 0.08, C['hkred'])

    add_text_box(slide, '3個簡單步驟', 0.5, 0.3, 6, 0.7,
                font_size=32, bold=True, color=C['dark'])
    add_text_box(slide, '至', 6.5, 0.4, 0.8, 0.6,
                font_size=24, color=C['muted'])
    add_text_box(slide, '3-5個工作天成功開店', 7, 0.3, 5.5, 0.7,
                font_size=28, bold=True, color=C['hkred'])

    steps = [
        ('開立帳戶', '商業登記證(BR)\n銀行月結單'),
        ('建立電子合約', '完成合約\n簽署程序'),
        ('開展您的業務', '開始上架產品\n並接受訂單')
    ]
    x = 1
    for name, desc in steps:
        # 步驟編號圓圈
        add_rectangle(slide, x, 1.5, 3.5, 2.5, C['hkred'])
        add_text_box(slide, name, x + 0.2, 1.7, 3.1, 0.6,
                    font_size=20, bold=True, color=C['white'], align=PP_ALIGN.CENTER)
        add_text_box(slide, desc, x + 0.2, 2.4, 3.1, 1.2,
                    font_size=14, color=C['white'], align=PP_ALIGN.CENTER)
        x += 4

    # ────────────────────────────────────────────────────────
    # SLIDE 17: 痛點解決方案（動態）
    # ────────────────────────────────────────────────────────
    if pain_points:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_background(slide, C['white'])
        add_rectangle(slide, 0, 0, 13.333, 0.08, C['hkred'])

        add_text_box(slide, 'HKTVmall 助您解決商業挑戰', 0.5, 0.3, 12.333, 0.7,
                    font_size=28, bold=True, color=C['dark'])

        y = 1.2
        for pain_key in pain_points:
            if pain_key in PAIN_SOLUTIONS:
                data = PAIN_SOLUTIONS[pain_key]

                # 痛點
                add_rectangle(slide, 0.5, y, 0.15, 0.5, C['hkred'])
                add_text_box(slide, f'❌ {data["zh"]}', 0.8, y, 5, 0.5,
                            font_size=16, bold=True, color=C['hkred'])

                # 解決方案
                add_rectangle(slide, 0.5, y + 0.5, 0.15, 0.5, C['green'])
                add_text_box(slide, f'✓ {data["solution"]}', 0.8, y + 0.5, 11.5, 0.5,
                            font_size=13, color=C['green'])

                y += 1.3

    # ────────────────────────────────────────────────────────
    # SLIDE 18: 聯絡我們
    # ────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, C['hkred'])

    add_rectangle(slide, 0, 0, 13.333, 0.15, C['gold'])
    add_rectangle(slide, 0, 7.35, 13.333, 0.15, C['gold'])

    add_text_box(slide, '立即加入我們', 0.5, 1.2, 12.333, 1.2,
                font_size=56, bold=True, color=C['white'], align=PP_ALIGN.CENTER)

    add_text_box(slide, f'電郵地址: {contact_email}',
                0.5, 3, 12.333, 0.6,
                font_size=24, color=C['white'], align=PP_ALIGN.CENTER)
    add_text_box(slide, '網站: http://business.hktvmall.com',
                0.5, 3.8, 12.333, 0.5,
                font_size=20, color=C['white'], align=PP_ALIGN.CENTER)
    add_text_box(slide, 'WHATSAPP: +852 5283 4138',
                0.5, 4.4, 12.333, 0.5,
                font_size=20, color=C['white'], align=PP_ALIGN.CENTER)
    add_text_box(slide, 'https://wa.me/85252834138',
                0.5, 5, 12.333, 0.4,
                font_size=16, color=C['gold'], align=PP_ALIGN.CENTER)
    add_text_box(slide, f'{merchant_name} 期待與您合作！',
                0.5, 6, 12.333, 0.5,
                font_size=22, bold=True, color=C['gold'], align=PP_ALIGN.CENTER)

    # ────────────────────────────────────────────────────────
    # SLIDE 19: 附錄
    # ────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, C['white'])
    add_rectangle(slide, 0, 0, 13.333, 0.08, C['hkred'])

    add_text_box(slide, '附錄', 0.5, 0.3, 12.333, 0.7,
                font_size=32, bold=True, color=C['dark'])

    appendix_points = [
        '年費續約安排',
        'HKTVmall 商戶專屬 GOGOVAN 物流優惠',
        '新商戶入帳 HKTVmall 銀行資料',
        '商戶管理平台功能概覽',
        '集運計劃營運模式'
    ]

    y = 1.5
    for point in appendix_points:
        add_text_box(slide, f'•  {point}', 0.8, y, 11, 0.5,
                    font_size=18, color=C['dark'])
        y += 0.6

    # 保存
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output


# ============================================================
# API 端點
# ============================================================

@app.route('/generate', methods=['POST'])
def generate():
    """PPT 生成 API"""
    try:
        data = request.get_json()

        merchant_name = data.get('merchant_name', '')
        category = data.get('category', '本地')
        pain_points = data.get('pain_points', [])
        contact_email = data.get('contact_email', 'aog.merc@hktv.com.hk')

        ppt_file = create_pptx(
            merchant_name=merchant_name,
            category=category,
            pain_points=pain_points,
            contact_email=contact_email
        )

        filename = f'HKTVmall_招商方案_{merchant_name}.pptx'

        return send_file(
            ppt_file,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name=filename
        )
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/health', methods=['GET'])
def health():
    """健康檢查"""
    return jsonify({'status': 'ok'})


@app.route('/')
def index():
    """首頁"""
    return '''
    <html>
    <head><title>HKTVmall PPT Generator API</title></head>
    <body>
    <h1>HKTVmall PPT 生成器 API</h1>
    <p>POST /generate with JSON body to generate PPT</p>
    <pre>
    {
        "merchant_name": "ABC電器",
        "category": "本地|海外|服務",
        "pain_points": ["logistics", "fee"],
        "contact_email": "aog.merc@hktv.com.hk"
    }
    </pre>
    </body>
    </html>
    '''


# ============================================================
# 主程序
# ============================================================

def main():
    parser = argparse.ArgumentParser(description='HKTVmall 商戶招商 PPT 生成器 API')
    parser.add_argument('--host', default='0.0.0.0', help='主機 (默認 0.0.0.0)')
    parser.add_argument('--port', type=int, default=5000, help='端口 (默認 5000)')
    parser.add_argument('--debug', action='store_true', help='調試模式')

    args = parser.parse_args()

    print(f'''
╔══════════════════════════════════════════════════════════════╗
║     HKTVmall 商戶招商 PPT 生成器 API                      ║
╠══════════════════════════════════════════════════════════════╣
║  API 地址: http://{args.host}:{args.port}                       ║
║  文檔:    http://{args.host}:{args.port}/                       ║
║                                                      ║
║  使用方法:                                             ║
║    POST /generate                                     ║
║    {{                                                   ║
║      "merchant_name": "ABC電器",                         ║
║      "category": "本地",                                ║
║      "pain_points": ["logistics", "fee"],                ║
║      "contact_email": "aog.merc@hktv.com.hk"             ║
║    }}                                                   ║
╚══════════════════════════════════════════════════════════════╝
    ''')

    app.run(host=args.host, port=args.port, debug=args.debug)


if __name__ == '__main__':
    main()
