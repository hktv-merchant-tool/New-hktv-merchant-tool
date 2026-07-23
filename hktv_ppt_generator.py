#!/usr/bin/env python3
"""
HKTVmall 商戶招商 PPT 生成器 v3
基於 Short Version 模板 (22頁)

使用方法:
  python hktv_ppt_generator.py                                    # 生成示例 PPT
  python hktv_ppt_generator.py --serve                           # 啟動 Flask API 服務
  python hktv_ppt_generator.py --merchant "ABC電器" --pains "物流成本高,平台費用不透明"

指令 1: 頁面 [1, 3-22] 保留，頁 1 加入商戶名稱
指令 2: 頁面 [2] 根據痛點選擇動態生成

API 端點: POST /generate
  JSON body: {
    "merchant_name": "ABC電器有限公司",
    "category": "本地" | "海外" | "服務",
    "pain_points": ["物流成本高", "平台費用不透明"],
    "contact_email": "aog.merc@hktv.com.hk"
  }
"""

import io
import os
import sys
import json
import argparse
from pathlib import Path

from flask import Flask, request, send_file, jsonify
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor as RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

# ============================================================
# 模板數據 - 基於 Short Version 模板
# ============================================================

# 痛點與解決方案
PAIN_SOLUTIONS = {
    "物流成本高": {
        "pain": "物流成本高",
        "solution": "全港最大住宅配送網絡，350+輛貨車，500+取貨點，大幅降低物流成本"
    },
    "平台費用不透明": {
        "pain": "平台費用不透明",
        "solution": "佣金按分類碼固定比例計算，每月15個工作日結算，無隱藏收費"
    },
    "缺乏電商經驗": {
        "pain": "缺乏電商經驗",
        "solution": "電商學院提供完整培訓，MMS後台操作、產品上架、廣告投放一對一支援"
    },
    "擔心庫存風險": {
        "pain": "擔心庫存風險",
        "solution": "GREEN LAB臨期百貨計劃，3PL靈活存貨管理，代售模式先售後買降低風險"
    },
    "希望提升品牌知名度": {
        "pain": "希望提升品牌知名度",
        "solution": "每年過億廣告投放，TVB黃金時段，58個MTR站，23列全車身，260萬+買家"
    },
    "希望拓展年輕客群": {
        "pain": "希望拓展年輕客群",
        "solution": "精準CRM定位，25-40歲核心消費群，每季度4.7x購買頻率的忠實客戶"
    },
    "希望增加線上曝光": {
        "pain": "希望增加線上曝光",
        "solution": "關鍵字廣告、橫幅廣告、首頁推廣位，SEO及站內搜尋排名優化"
    }
}

# 招商計劃詳情
PLANS = {
    "本地": {
        "name": "常規商戶加盟計劃",
        "annual_fee": "HK$25,000 + HK$22,000 廣告金",
        "description": "本地常規商戶"
    },
    "海外": {
        "name": "Overseas Plan - 海外商戶",
        "annual_fee": "HK$3,000（原 HK$25,000）",
        "description": "海外跨境商戶"
    },
    "服務": {
        "name": "Service Deal - 服務類商戶",
        "annual_fee": "HK$3,000",
        "description": "服務類商戶"
    }
}

# ============================================================
# PPT 幻燈片內容 - 22頁 Short Version 結構
# ============================================================

def create_slide(prs):
    """創建空白幻燈片"""
    return prs.slides.add_slide(prs.slide_layouts[6])

def add_text(slide, text, x, y, w, h, size, bold=False, color=None, align=PP_ALIGN.LEFT):
    """添加文字框"""
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    if color:
        p.font.color.rgb = RgbColor(*color)
    p.alignment = align
    return shape

# ─── 第1頁: 封面 ───
def create_slide_01(prs, merchant_name=None):
    slide = create_slide(prs)
    add_text(slide, "No.1", 0.5, 1.5, 9, 0.9, 56, bold=True, color=(230, 0, 18), align=PP_ALIGN.CENTER)
    add_text(slide, "香港網上購物平台", 0.5, 2.4, 9, 0.8, 40, bold=True, color=(34, 34, 34), align=PP_ALIGN.CENTER)
    add_text(slide, "商戶合作方案", 0.5, 3.3, 9, 0.7, 36, color=(34, 34, 34), align=PP_ALIGN.CENTER)
    if merchant_name:
        add_text(slide, f"「{merchant_name}」專屬方案", 0.5, 4.5, 9, 0.6, 24, bold=True, color=(230, 0, 18), align=PP_ALIGN.CENTER)

# ─── 第2頁: 商戶資料【主要業務】和【痛點與解決方案】───
def create_slide_02(prs, pain_points):
    slide = create_slide(prs)
    add_text(slide, "商戶資料", 0.5, 0.3, 9, 0.6, 26, bold=True, color=(34, 34, 34), align=PP_ALIGN.CENTER)
    add_text(slide, "【痛點與解決方案】", 0.8, 1.0, 8.5, 0.5, 16, bold=True, color=(0, 102, 204))

    y = 1.5
    for pain_key in pain_points:
        if pain_key in PAIN_SOLUTIONS:
            data = PAIN_SOLUTIONS[pain_key]
            add_text(slide, f"❌ {data['pain']}", 0.8, y, 4, 0.35, 12, bold=True, color=(204, 0, 0))
            add_text(slide, f"✓ {data['solution']}", 4.8, y, 4.7, 0.35, 10, color=(0, 128, 0))
            y += 0.55

# ─── 第3頁: 發展海外市場 ───
def create_slide_03(prs):
    slide = create_slide(prs)
    add_text(slide, "發展海外市場", 0.5, 0.3, 9, 0.6, 26, bold=True, color=(34, 34, 34))
    items = [
        ("營銷與廣告", "過億重磅廣告投放，TVC電視廣告"),
        ("物流 & O2O 門市", "75家O2O門店，500+取貨點"),
        ("24小時網上購物", "全天候便捷購物體驗"),
        ("平台大數據分析", "260萬+活躍買家，精準CRM定位"),
        ("商客互聯平台與直播頻道", "直播帶貨，實時互動提升成交")
    ]
    y = 1.1
    for title, desc in items:
        add_text(slide, title, 0.8, y, 4, 0.4, 14, bold=True, color=(230, 0, 18))
        add_text(slide, desc, 0.8, y + 0.35, 8, 0.4, 12, color=(34, 34, 34))
        y += 0.85

# ─── 第4頁: HKTVmall 的里程碑 ───
def create_slide_04(prs):
    slide = create_slide(prs)
    add_text(slide, "HKTVmall 的里程碑", 0.5, 0.25, 9, 0.5, 24, bold=True, color=(34, 34, 34))
    add_text(slide, "城市電訊在香港聯合交易所上市（香港聯交所代號：1137）", 0.5, 0.7, 9, 0.3, 10, color=(128, 128, 128), align=PP_ALIGN.CENTER)
    milestones = [
        ("1992", "城市電訊（香港）有限公司在香港正式成立"),
        ("1997", "香港電視面世製作電視節目"),
        ("2000", "附屬公司香港寬頻推出網路服務"),
        ("2012", "電商正式啟動"),
        ("2015–2016", "香港電視正式轉型為網上購物"),
        ("2017–2018", "物流中心自動化揀貨及倉存系統全面運作"),
        ("2019–2020", "海外業務擴張及開拓業務新板塊"),
        ("2021–2022", "街市即日餸、Mall Dollar回贈"),
        ("2023–現在", "持續引領香港電商市場")
    ]
    y = 1.05
    for year, desc in milestones:
        add_text(slide, year, 0.6, y, 1.3, 0.4, 10, bold=True, color=(230, 0, 18))
        add_text(slide, desc, 2.0, y, 7, 0.4, 10, color=(34, 34, 34))
        y += 0.46

# ─── 第5頁: 強大流量與卓越表現 ───
def create_slide_05(prs):
    slide = create_slide(prs)
    add_text(slide, "強大流量與卓越表現", 0.5, 0.3, 9, 0.6, 26, bold=True, color=(34, 34, 34))
    add_text(slide, "香港領先的電子商務平台", 0.5, 1.0, 9, 0.4, 16, color=(128, 128, 128), align=PP_ALIGN.CENTER)
    add_text(slide, "2025 Total GMV", 0.5, 1.8, 9, 0.5, 18, color=(128, 128, 128), align=PP_ALIGN.CENTER)
    add_text(slide, "HK$7.89 Billion", 0.5, 2.3, 9, 1.2, 52, bold=True, color=(230, 0, 18), align=PP_ALIGN.CENTER)
    add_text(slide, "*香港科技探索有限公司2025年年報數據", 0.5, 3.6, 9, 0.3, 10, color=(128, 128, 128), align=PP_ALIGN.CENTER)

# ─── 第6頁: 全方位的產品與服務 ───
def create_slide_06(prs):
    slide = create_slide(prs)
    add_text(slide, "全方位的產品與服務", 0.5, 0.3, 9, 0.6, 26, bold=True, color=(34, 34, 34))

# ─── 第7頁: 2025年各類別的銷售表現 ───
def create_slide_07(prs):
    slide = create_slide(prs)
    add_text(slide, "2025年各類別的銷售表現", 0.5, 0.3, 9, 0.6, 26, bold=True, color=(34, 34, 34))
    add_text(slide, "2025 Total GMV", 0.5, 1.3, 4, 0.4, 14, color=(128, 128, 128))
    add_text(slide, "HK$7.89", 0.5, 1.7, 4, 1, 44, bold=True, color=(230, 0, 18))
    add_text(slide, "Billion", 3.5, 2.1, 2, 0.5, 18, color=(230, 0, 18))
    add_text(slide, "*香港科技探索有限公司2025年年報數據", 0.5, 2.8, 9, 0.3, 10, color=(128, 128, 128))

# ─── 第8頁: 有利商戶的網購平台優勢 ───
def create_slide_08(prs):
    slide = create_slide(prs)
    add_text(slide, "有利商戶的網購平台優勢", 0.5, 0.3, 9, 0.6, 26, bold=True, color=(34, 34, 34))

# ─── 第9頁: 國際知名品牌 ───
def create_slide_09(prs):
    slide = create_slide(prs)
    add_text(slide, "國際知名品牌遍布每個產品類別", 0.5, 0.3, 9, 0.6, 26, bold=True, color=(34, 34, 34))

# ─── 第10頁: 全面滲透，擁有高回購率 ───
def create_slide_10(prs):
    slide = create_slide(prs)
    add_text(slide, "全面滲透，擁有高回購率", 0.5, 0.3, 9, 0.6, 26, bold=True, color=(34, 34, 34))
    items = ["客戶性別分佈均衡", "客戶橫誇各個年齡層", "均衡的用戶基礎助力商家精準觸達男女消費群體", "每季度4.7x購買 = 忠實客戶", "高重覆購買率 = 為商家提供穩定的銷售"]
    y = 1.1
    for item in items:
        add_text(slide, f"• {item}", 1, y, 8, 0.5, 14, color=(34, 34, 34))
        y += 0.55

# ─── 第11頁: 五大地區，七大貨倉 ───
def create_slide_11(prs):
    slide = create_slide(prs)
    add_text(slide, "五大地區，七大貨倉", 0.5, 0.3, 9, 0.6, 26, bold=True, color=(34, 34, 34))
    warehouses = [("上水", "常温儲運倉"), ("屯門", "常温儲運倉"), ("青衣", "常温儲運倉"), ("葵涌", "3PL物流倉庫"), ("將軍澳", "倉庫總部貨倉")]
    y = 1.1
    for name, desc in warehouses:
        add_text(slide, name, 0.7, y, 1.5, 0.55, 15, bold=True, color=(230, 0, 18))
        add_text(slide, desc, 2.4, y, 6, 0.55, 14, color=(34, 34, 34))
        y += 0.62
    add_text(slide, "• 常温/冷凍/急凍儲運倉", 0.7, 4.3, 8, 0.4, 12, color=(128, 128, 128))

# ─── 第12頁: 全港最大住宅配送網絡 ───
def create_slide_12(prs):
    slide = create_slide(prs)
    add_text(slide, "全港最大住宅配送網絡", 0.5, 0.3, 9, 0.6, 26, bold=True, color=(34, 34, 34))
    add_text(slide, "500+", 0.5, 1.2, 9, 1.3, 80, bold=True, color=(230, 0, 18), align=PP_ALIGN.CENTER)
    add_text(slide, "取貨點", 0.5, 2.6, 9, 0.6, 26, color=(34, 34, 34), align=PP_ALIGN.CENTER)
    add_text(slide, "(香港郵政、Circle K 便利店、759阿信屋，真好城等)", 0.5, 3.3, 9, 0.4, 11, color=(128, 128, 128), align=PP_ALIGN.CENTER)

# ─── 第13頁: 大型優惠推廣活動 ───
def create_slide_13(prs):
    slide = create_slide(prs)
    add_text(slide, "大型優惠推廣活動", 0.5, 0.3, 9, 0.6, 26, bold=True, color=(34, 34, 34))
    add_text(slide, "單日銷售額突破", 0.5, 1.0, 9, 0.5, 18, color=(34, 34, 34), align=PP_ALIGN.CENTER)
    add_text(slide, "4,000萬", 0.5, 1.5, 9, 1.1, 56, bold=True, color=(230, 0, 18), align=PP_ALIGN.CENTER)
    add_text(slide, "HKTVmall全額承擔85折", 0.5, 2.8, 9, 0.8, 14, color=(128, 128, 128), align=PP_ALIGN.CENTER)

# ─── 第14頁: 全年不間斷為客戶的推廣活動 ───
def create_slide_14(prs):
    slide = create_slide(prs)
    add_text(slide, "全年不間斷為客戶的推廣活動", 0.5, 0.3, 9, 0.6, 26, bold=True, color=(34, 34, 34))

# ─── 第15頁: 顧客互動工具 ───
def create_slide_15(prs):
    slide = create_slide(prs)
    add_text(slide, "顧客互動工具", 0.5, 0.3, 9, 0.6, 26, bold=True, color=(34, 34, 34))
    add_text(slide, "有效提升銷售並建立長期忠誠度", 0.5, 1.0, 9, 0.4, 16, color=(128, 128, 128), align=PP_ALIGN.CENTER)
    items = ["直播限定優惠為消費者帶來購物驚喜", "連結商家與消費者進行直接互動，提升即時成交機會"]
    y = 1.8
    for item in items:
        add_text(slide, f"• {item}", 1, y, 8, 0.6, 14, color=(34, 34, 34))
        y += 0.65

# ─── 第16頁: 多元化互動工具 ───
def create_slide_16(prs):
    slide = create_slide(prs)
    add_text(slide, "多元化互動工具", 0.5, 0.3, 9, 0.6, 22, bold=True, color=(34, 34, 34))
    add_text(slide, "提升客戶使用率", 3.5, 0.35, 5, 0.5, 14, color=(128, 128, 128))
    add_text(slide, "ChicChat", 0.7, 1.0, 8, 0.4, 16, bold=True, color=(230, 0, 18))
    add_text(slide, "「查詢商家」功能", 0.7, 1.4, 4, 0.4, 13, bold=True, color=(34, 34, 34))
    add_text(slide, "透過即時對話功能提升服務效率", 0.7, 1.75, 8, 0.4, 11, color=(128, 128, 128))
    add_text(slide, "客戶群組系統", 0.7, 2.3, 4, 0.4, 13, bold=True, color=(34, 34, 34))
    items = ["允許商家、消費者和KOL交流互動", "允許商家與目標客戶討論新事物、新話題", "一站式自助工具，用以精準劃分客戶群"]
    y = 2.7
    for item in items:
        add_text(slide, item, 0.7, y, 8, 0.4, 11, color=(128, 128, 128))
        y += 0.4

# ─── 第17頁: 商戶數據分析平台 ───
def create_slide_17(prs):
    slide = create_slide(prs)
    add_text(slide, "商戶數據分析平台", 0.5, 0.3, 9, 0.6, 22, bold=True, color=(34, 34, 34))
    add_text(slide, "輕鬆剖析業務並進一步策劃營銷", 0.5, 0.85, 8, 0.4, 14, color=(128, 128, 128))
    items = [
        ("銷售數據與消費者分析", "銷售額總覽、產品和訂單趨勢"),
        ("分析推廣活動及廣告效益", "設定價格和促銷機制"),
        ("客戶人口分群及統計資料", "評估自身促銷活動的成效"),
        ("內外流量數據", "")
    ]
    y = 1.4
    for title, desc in items:
        add_text(slide, title, 0.7, y, 8, 0.4, 13, bold=True, color=(34, 34, 34))
        if desc:
            add_text(slide, f"• {desc}", 0.9, y + 0.3, 8, 0.35, 11, color=(128, 128, 128))
        y += 0.7

# ─── 第18頁: 簡易加入流程 ───
def create_slide_18(prs):
    slide = create_slide(prs)
    add_text(slide, "3個簡單步驟", 0.5, 0.3, 4.5, 0.6, 24, bold=True, color=(34, 34, 34))
    add_text(slide, "至", 4.5, 0.35, 0.5, 0.5, 18, color=(128, 128, 128))
    add_text(slide, "3至5個工作天成功開店", 5, 0.3, 4.5, 0.6, 24, bold=True, color=(230, 0, 18))
    steps = [("開立帳戶", "商業登記證(BR)、銀行月結單"), ("建立電子合約", "完成合約簽署程序"), ("開展您的業務", "開始上架產品並接受訂單")]
    y = 1.2
    for name, desc in steps:
        add_text(slide, name, 0.7, y, 3, 0.45, 16, bold=True, color=(230, 0, 18))
        add_text(slide, desc, 4, y, 5, 0.45, 12, color=(34, 34, 34))
        y += 0.75

# ─── 第19頁: 常規商戶加盟計劃 ───
def create_slide_19(prs, plan_key="本地"):
    plan = PLANS.get(plan_key, PLANS["本地"])
    slide = create_slide(prs)
    add_text(slide, plan["name"], 0.5, 0.3, 9, 0.6, 24, bold=True, color=(34, 34, 34))
    if plan_key == "本地":
        add_text(slide, "HK$25,000", 0.5, 1.2, 4, 0.8, 36, bold=True, color=(230, 0, 18))
        add_text(slide, "+", 4.2, 1.2, 0.6, 0.8, 28, color=(34, 34, 34), align=PP_ALIGN.CENTER)
        add_text(slide, "HKTVmall廣告金", 4.8, 1.2, 4, 0.8, 18, color=(34, 34, 34))
    else:
        add_text(slide, plan["annual_fee"], 0.5, 1.2, 9, 0.8, 32, bold=True, color=(230, 0, 18))
    add_text(slide, "*廣告金回贈按比例發放，並需於當年12月31日前用完", 0.5, 2.6, 9, 0.4, 11, color=(128, 128, 128))

# ─── 第20頁: 數碼引流計劃 ───
def create_slide_20(prs):
    slide = create_slide(prs)
    add_text(slide, "HKTVmall數碼引流計劃", 0.5, 0.3, 9, 0.6, 22, bold=True, color=(34, 34, 34))
    add_text(slide, "尊享特惠佣金率*", 0.5, 1.2, 9, 0.6, 20, bold=True, color=(0, 136, 0))
    add_text(slide, "於其他社交媒體推廣/投放廣告", 0.5, 2.0, 9, 0.5, 16, color=(34, 34, 34))
    add_text(slide, "將連結引流到 HKTVmall", 0.5, 2.5, 9, 0.5, 16, color=(34, 34, 34))

# ─── 第21頁: 貨款每月結算流程 ───
def create_slide_21(prs):
    slide = create_slide(prs)
    add_text(slide, "貨款每月結算流程", 0.5, 0.3, 9, 0.6, 24, bold=True, color=(34, 34, 34))
    add_text(slide, "本月", 0.7, 1.2, 2, 0.45, 16, bold=True, color=(230, 0, 18))
    add_text(slide, "本月完成的訂單", 2.7, 1.2, 6, 0.45, 14, color=(34, 34, 34))
    add_text(slide, "在15個工作日內", 0.7, 1.8, 3, 0.45, 16, bold=True, color=(230, 0, 18))
    add_text(slide, "款項轉至商戶的銀行帳戶", 3.7, 1.8, 5, 0.45, 14, color=(34, 34, 34))

# ─── 第22頁: 聯絡我們 ───
def create_slide_22(prs, merchant_name=None):
    slide = create_slide(prs)
    add_text(slide, "立即加入我們", 0.5, 1.0, 9, 0.9, 36, bold=True, color=(34, 34, 34), align=PP_ALIGN.CENTER)
    add_text(slide, "電郵地址 : aog.merc@hktv.com.hk", 0.5, 2.2, 9, 0.5, 20, align=PP_ALIGN.CENTER)
    add_text(slide, "網站 : http://business.hktvmall.com", 0.5, 2.8, 9, 0.5, 18, align=PP_ALIGN.CENTER)
    add_text(slide, "WHATSAPP : +852 5283 4138", 0.5, 3.4, 9, 0.5, 18, align=PP_ALIGN.CENTER)
    add_text(slide, "https://wa.me/85252834138", 0.5, 3.9, 9, 0.4, 14, color=(128, 128, 128), align=PP_ALIGN.CENTER)
    if merchant_name:
        add_text(slide, f"「{merchant_name}」期待與您合作！", 0.5, 4.6, 9, 0.5, 20, bold=True, color=(230, 0, 18), align=PP_ALIGN.CENTER)

# ============================================================
# 主 PPT 生成函數
# ============================================================

def generate_merchant_ppt(merchant_name=None, category="本地", pain_points=None, contact_email="aog.merc@hktv.com.hk", output_path=None):
    """
    生成商戶招商 PPT - Short Version (22頁)

    指令 1: 頁面 [1, 3-22] 保留，頁 1 加入商戶名稱
    指令 2: 頁面 [2] 根據痛點選擇動態生成
    """
    if pain_points is None:
        pain_points = []

    # 創建演示文稿（10:5.625 = 16:9）
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # 第1頁: 封面 + 商戶名稱
    create_slide_01(prs, merchant_name)

    # 第2頁: 痛點與解決方案（根據選擇動態生成）
    if pain_points:
        create_slide_02(prs, pain_points)

    # 第3頁: 發展海外市場
    create_slide_03(prs)

    # 第4頁: HKTVmall 的里程碑
    create_slide_04(prs)

    # 第5頁: 強大流量與卓越表現
    create_slide_05(prs)

    # 第6頁: 全方位的產品與服務
    create_slide_06(prs)

    # 第7頁: 2025年各類別的銷售表現
    create_slide_07(prs)

    # 第8頁: 有利商戶的網購平台優勢
    create_slide_08(prs)

    # 第9頁: 國際知名品牌
    create_slide_09(prs)

    # 第10頁: 全面滲透，擁有高回購率
    create_slide_10(prs)

    # 第11頁: 五大地區，七大貨倉
    create_slide_11(prs)

    # 第12頁: 全港最大住宅配送網絡
    create_slide_12(prs)

    # 第13頁: 大型優惠推廣活動
    create_slide_13(prs)

    # 第14頁: 全年不間斷為客戶的推廣活動
    create_slide_14(prs)

    # 第15頁: 顧客互動工具
    create_slide_15(prs)

    # 第16頁: 多元化互動工具
    create_slide_16(prs)

    # 第17頁: 商戶數據分析平台
    create_slide_17(prs)

    # 第18頁: 簡易加入流程
    create_slide_18(prs)

    # 第19頁: 常規商戶加盟計劃
    create_slide_19(prs, category)

    # 第20頁: 數碼引流計劃
    create_slide_20(prs)

    # 第21頁: 貨款每月結算流程
    create_slide_21(prs)

    # 第22頁: 聯絡我們
    create_slide_22(prs, merchant_name)

    # 保存
    if output_path:
        prs.save(output_path)
        print(f"PPT 已保存至: {output_path}")
    else:
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        return output

# ============================================================
# Flask API
# ============================================================

app = Flask(__name__)

@app.route('/generate', methods=['POST'])
def generate():
    """PPT 生成 API"""
    try:
        data = request.get_json()
        merchant_name = data.get('merchant_name', '')
        category = data.get('category', '本地')
        pain_points = data.get('pain_points', [])
        contact_email = data.get('contact_email', 'aog.merc@hktv.com.hk')

        ppt_file = generate_merchant_ppt(
            merchant_name=merchant_name,
            category=category,
            pain_points=pain_points,
            contact_email=contact_email
        )

        return send_file(
            ppt_file,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name=f'HKTVmall_招商方案_{merchant_name}.pptx'
        )
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/health', methods=['GET'])
def health():
    return jsonify({'status': 'ok'})

# ============================================================
# 主程序
# ============================================================

def main():
    parser = argparse.ArgumentParser(description='HKTVmall 商戶招商 PPT 生成器 v3 (Short Version 22頁)')
    parser.add_argument('--serve', action='store_true', help='啟動 Flask API 服務')
    parser.add_argument('--output', '-o', help='輸出文件路徑')
    parser.add_argument('--merchant', '-m', help='商戶名稱')
    parser.add_argument('--category', '-c', choices=['本地', '海外', '服務'], default='本地', help='招商計劃')
    parser.add_argument('--pains', '-p', help='痛點（逗號分隔）')
    parser.add_argument('--email', '-e', default='aog.merc@hktv.com.hk', help='聯絡電郵')
    parser.add_argument('--host', default='0.0.0.0', help='API 主機')
    parser.add_argument('--port', type=int, default=5000, help='API 端口')

    args = parser.parse_args()

    if args.serve:
        print(f"啟動 API 服務: http://{args.host}:{args.port}")
        app.run(host=args.host, port=args.port, debug=True)
    else:
        pain_points = []
        if args.pains:
            pain_points = [p.strip() for p in args.pains.split(',')]

        output_path = args.output or f'HKTVmall_招商方案_{args.merchant or "示例"}.pptx'
        print(f"生成 PPT: {output_path}")
        print(f"  商戶: {args.merchant or '示例'}")
        print(f"  痛點: {pain_points or '無'}")

        generate_merchant_ppt(
            merchant_name=args.merchant,
            category=args.category,
            pain_points=pain_points,
            contact_email=args.email,
            output_path=output_path
        )

if __name__ == '__main__':
    main()
